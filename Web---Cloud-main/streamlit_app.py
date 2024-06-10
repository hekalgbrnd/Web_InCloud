import streamlit as st
import os
import json
from pathlib import Path
import shutil
from streamlit_option_menu import option_menu
from docx import Document
import openpyxl
from pptx import Presentation
import pandas as pd

# Set Streamlit page configuration
st.set_page_config(page_title="In - Cloud", layout="wide", initial_sidebar_state="expanded")

# Define function to save the favorite state
def save_favorites():
    with open("favorites.json", "w") as f:
        json.dump({"folders": folder_favorites, "files": file_favorites}, f)

# Load or initialize the favorite state
FAVORITES_FILE = "favorites.json"
if Path(FAVORITES_FILE).exists():
    try:
        with open(FAVORITES_FILE, "r") as f:
            favorites_data = json.load(f)
            if isinstance(favorites_data, dict):  # Ensure it's a dictionary
                folder_favorites = favorites_data.get("folders", [])
                file_favorites = favorites_data.get("files", [])
            else:
                folder_favorites = []
                file_favorites = []
                save_favorites()  # Re-initialize favorites file with correct structure
    except json.JSONDecodeError:
        folder_favorites = []
        file_favorites = []
        save_favorites()  # Re-initialize favorites file with correct structure
else:
    folder_favorites = []
    file_favorites = []

# Check if a folder or file is a favorite
def is_favorite(item, is_folder=True):
    return str(item) in (folder_favorites if is_folder else file_favorites)

# Toggle favorite status of a folder or file
def toggle_favorite(item, is_folder=True):
    item_str = str(item)
    if is_folder:
        if item_str in folder_favorites:
            folder_favorites.remove(item_str)
        else:
            folder_favorites.append(item_str)
    else:
        if item_str in file_favorites:
            file_favorites.remove(item_str)
        else:
            file_favorites.append(item_str)
    save_favorites()

# Add custom CSS for styling
st.markdown("""
    <style>
    body {
        background-image: linear-gradient(to bottom right, #4CAF50, #8bc34a);
        color: #fff;
    }
    .button-icon {
        display: flex;
        align-items: center;
    }
    .button-icon i {
        margin-right: 5px;
    }
    </style>
""", unsafe_allow_html=True)

with st.sidebar:
    selected = option_menu("InClouds", ["About Us", "Home", "Favorite"],
                           icons=["info", "house", "heart"], menu_icon="cast", default_index=1,
                           styles={"nav-link-selected": {"background-color": "#8bc34a"}})

# Configure main directory
BASE_DIR = Path("uploads")
if not BASE_DIR.exists():
    BASE_DIR.mkdir(parents=True, exist_ok=True)

# Function to create a folder if it does not exist
def create_folder(path):
    if not path.exists():
        path.mkdir(parents=True, exist_ok=True)

# Function to add a document to a folder
def add_document(folder, file):
    try:
        folder_path = BASE_DIR / folder
        create_folder(folder_path)
        file_path = folder_path / file.name
        with open(file_path, "wb") as f:
            f.write(file.getbuffer())
        st.success(f"File '{file.name}' successfully uploaded to folder '{folder}'.")
    except Exception as e:
        st.error(f"Failed to add document: {e}")

# Function to rename a file
def rename_file(file_path, new_name):
    try:
        new_file_path = file_path.parent / new_name
        file_path.rename(new_file_path)
        st.success(f"File '{file_path.name}' successfully renamed to '{new_name}'.")
    except Exception as e:
        st.error(f"Failed to rename file: {e}")

# Function to delete a file
def delete_file(file_path):
    try:
        file_path.unlink()
        st.success(f"File '{file_path.name}' successfully deleted.")
    except Exception as e:
        st.error(f"Failed to delete file: {e}")

# Function to delete a folder
def delete_folder(folder_path):
    try:
        shutil.rmtree(folder_path)
        st.success(f"Folder '{folder_path.name}' successfully deleted along with its contents.")
    except Exception as e:
        st.error(f"Failed to delete folder: {e}")

# Function to rename a folder
def rename_folder(folder_path, new_name):
    try:
        new_folder_path = folder_path.parent / new_name
        folder_path.rename(new_folder_path)
        st.success(f"Folder '{folder_path.name}' successfully renamed to '{new_name}'.")
    except Exception as e:
        st.error(f"Failed to rename folder: {e}")

# Function to calculate total storage size
def calculate_total_size(path):
    total_size = 0
    for dirpath, dirnames, filenames in os.walk(path):
        for f in filenames:
            fp = os.path.join(dirpath, f)
            total_size += os.path.getsize(fp)
    return total_size

# Function to display Word document content
def display_word_document(file_path):
    doc = Document(file_path)
    for paragraph in doc.paragraphs:
        st.write(paragraph.text)

# Function to display Excel document content
def display_excel_document(file_path):
    wb = openpyxl.load_workbook(file_path)
    sheet = wb.active
    data = sheet.values
    cols = next(data)
    data = list(data)
    df = pd.DataFrame(data, columns=cols)
    st.dataframe(df)

# Function to display PowerPoint presentation content
def display_ppt_document(file_path):
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    st.write(paragraph.text)

# Function to display folder contents
def show_folder_content(path, search_query=""):
    items = list(path.iterdir())
    files = [item for item in items if item.is_file() and search_query.lower() in item.name.lower()]
    folders = [item for item in items if item.is_dir() and search_query.lower() in item.name.lower()]

    st.write(f"Folder contents: {path.relative_to(BASE_DIR)}")
    if path != BASE_DIR and st.button("Back to main folder"):
        st.experimental_set_query_params(path="")

    st.write("### Folders")
    for folder in folders:
        col1, col2, col3, col4 = st.columns([5, 1, 1, 1])
        with col1:
            if st.button(f"📁 {folder.name}", key=f"folder_{folder.name}"):
                st.experimental_set_query_params(path=str(folder.relative_to(BASE_DIR)))
        with col2:
            if st.button("⭐" if is_favorite(folder, is_folder=True) else "☆", key=f"favorite_folder_{folder.name}"):
                toggle_favorite(folder, is_folder=True)
                st.experimental_rerun()
        with col3:
            if st.button("⚙️", key=f"settings_folder_{folder.name}"):
                st.session_state[f"menu_folder_{folder.name}"] = not st.session_state.get(f"menu_folder_{folder.name}", False)
            if st.session_state.get(f"menu_folder_{folder.name}", False):
                menu_options = ["Rename", "Delete"]
                action = st.selectbox("Select Action", menu_options, key=f"menu_action_{folder.name}")
                if action == "Rename":
                    new_name = st.text_input(f"Rename folder '{folder.name}'", key=f"rename_{folder.name}_input")
                    if st.button(f"Rename", key=f"btn_rename_{folder.name}_confirm"):
                        if new_name:
                            rename_folder(folder, new_name)
                            st.experimental_rerun()
                        else:
                            st.error("New folder name cannot be empty.")
                elif action == "Delete":
                    if st.button(f"Delete Folder", key=f"btn_delete_{folder.name}_confirm"):
                        delete_folder(folder)
                        st.experimental_rerun()
        with col4:
            st.write("")

    st.write("### Files")
    for file in files:
        col1, col2, col3, col4 = st.columns([5, 1, 1, 1])
        with col1:
            st.write(f"{file.name}")
            mime_type = None
            if file.suffix == ".docx":
                mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                if st.button("Display", key=f"display_{file.name}"):
                    display_word_document(file)
            elif file.suffix == ".xlsx":
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                if st.button("Display", key=f"display_{file.name}"):
                    display_excel_document(file)
            elif file.suffix == ".pptx":
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                if st.button("Display", key=f"display_{file.name}"):
                    display_ppt_document(file)

            if mime_type:
                google_url = upload_to_google_drive(file, mime_type)
                st.markdown(f"File can be accessed on Google Drive: [Open in Google Drive]({google_url})")
        
        with col2:
            if st.button("⭐" if is_favorite(file, is_folder=False) else "☆", key=f"favorite_file_{file.name}"):
                toggle_favorite(file, is_folder=False)
                st.experimental_rerun()
        
        with col3:
            if st.button("⚙️", key=f"settings_file_{file.name}"):
                st.session_state[f"menu_file_{file.name}"] = not st.session_state.get(f"menu_file_{file.name}", False)
            if st.session_state.get(f"menu_file_{file.name}", False):
                menu_options = ["Rename", "Delete", "Download"]
                action = st.selectbox("Select Action", menu_options, key=f"menu_action_{file.name}")
                if action == "Rename":
                    new_name = st.text_input(f"Rename file '{file.name}'", key=f"rename_{file.name}_input")
                    if st.button(f"Rename", key=f"btn_rename_{file.name}_confirm"):
                        if new_name:
                            rename_file(file, new_name)
                            st.experimental_rerun()
                        else:
                            st.error("New file name cannot be empty.")
                elif action == "Delete":
                    if st.button(f"Delete File", key=f"btn_delete_{file.name}_confirm"):
                        delete_file(file)
                        st.experimental_rerun()
                elif action == "Download":
                    with open(file, "rb") as f:
                        st.download_button(
                            label="Download",
                            data=f,
                            file_name=file.name,
                            mime="application/octet-stream",
                            key=f"download_{file.name}"
                        )
        with col4:
            st.write("")

# Streamlit Application
st.markdown("<h1 style='text-align: center; color: #fff; font-family: Arial, sans-serif; font-size: 50px;'>InClouds</h1>", unsafe_allow_html=True)

if selected == "Home":
    # Read query parameter for folder navigation
    query_params = st.experimental_get_query_params()
    current_path_str = query_params.get("path", [""])[0]
    current_path = BASE_DIR / current_path_str

    # Display total storage size
    total_storage_used = calculate_total_size(BASE_DIR)
    st.write(f"Total storage used: {total_storage_used / (1024 * 1024):.2f} MB")

    # Create a new folder
    st.header("Create Folder")
    folder_name = st.text_input("Folder Name")
    if st.button("Create Folder", key="create_folder"):
        if folder_name:
            create_folder(current_path / folder_name)
            st.success(f"Folder '{folder_name}' successfully created.")
        else:
            st.error("Folder name cannot be empty")

    # Upload a new file
    st.header("Upload File")
    uploaded_file = st.file_uploader("Choose File")
    if uploaded_file and st.button("Upload File", key="upload_file"):
        if uploaded_file.size <= 1 * 1024 * 1024 * 1024:  # Check if file size is less than or equal to 1GB
            add_document(current_path_str, uploaded_file)
        else:
            st.error("File size cannot exceed 1GB")

    # Search bar for searching files and folders
    st.header("Search")
    search_query = st.text_input("Search folder or file")
    search_button_clicked = st.button("Search", key="search_button")
    
    # Display current folder contents with search query
    show_folder_content(current_path, search_query)

elif selected == "Settings":
    st.write("Settings page content can go here.")

elif selected == "About Us":
    st.markdown("<h1 style='text-align: center; color: #fff; font-family: Arial, sans-serif; font-size: 30px;'>About Us</h1>", unsafe_allow_html=True)
    st.write("""
    Welcome to the About Us page!

    This Streamlit app is developed by YARSI Information Technology Faculty Cloud Computing team.
    """)

elif selected == "Favorite":
    st.title("Favorite Folders and Files")
   
    st.write("### Favorite Folders")
    for folder in folder_favorites:
        folder_path = Path(folder)
        if folder_path.is_dir() and st.button(f"📁 {folder_path.name}", key=f"favorite_folder_{folder_path.name}"):
            st.experimental_set_query_params(path=str(folder_path.relative_to(BASE_DIR)))

    st.write("### Favorite Files")
    for file in file_favorites:
        file_path = Path(file)
        if file_path.is_file():
            col1, col2, col3 = st.columns([6, 1, 1])
            with col1:
                st.write(f"{file_path.name}")
                if st.button("Display", key=f"display_{file_path.name}"):
                    if file_path.suffix == ".docx":
                        display_word_document(file_path)
                    elif file_path.suffix == ".xlsx":
                        display_excel_document(file_path)
                    elif file_path.suffix == ".pptx":
                        display_ppt_document(file_path)
            with col2:
                if st.button("Download", key=f"download_{file_path.name}"):
                    with open(file_path, "rb") as f:
                        st.download_button(
                            label="Download",
                            data=f,
                            file_name=file_path.name,
                            mime="application/octet-stream",
                            key=f"download_btn_{file_path.name}"
                        )
            with col3:
                if st.button("Remove from Favorites", key=f"remove_favorite_{file_path.name}"):
                    toggle_favorite(file_path, is_folder=False)
                    st.experimental_rerun()
